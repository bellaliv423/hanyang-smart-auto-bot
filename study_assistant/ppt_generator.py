"""
PPT 자동 생성기 - 수업 주제 → 전문적인 PPT 파일

기능:
1. 주제 입력 → Claude AI가 구조화된 PPT 내용 생성
2. 논문/기사/사례 검색 결과 포함
3. 한국어/中文/English 다국어 지원
4. Google Drive 자동 업로드
5. WhatsApp으로 완성 알림

사용법:
    python study_assistant/ppt_generator.py --topic "회귀분석" --course 196594
    python study_assistant/ppt_generator.py --topic "M&A 절차" --course 196656 --lang zh
    python study_assistant/ppt_generator.py --topic "계약법 기초" --slides 10 --lang ko
"""

import anthropic
import argparse
import io
import json
import os
import re
import sys
from datetime import datetime
from pathlib import Path

try:
    if sys.stdout.encoding != 'utf-8':
        sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
        sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')
except Exception:
    pass

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor

PROJECT_ROOT = Path(__file__).parent.parent

from dotenv import load_dotenv
load_dotenv(PROJECT_ROOT / "config" / ".env")

API_KEY = os.getenv("ANTHROPIC_API_KEY", "")

# 과목 매핑
COURSE_MAP = {
    196594: {"ko": "경영통계학", "en": "Business Statistics", "zh": "經營統計學", "prof": "부제만"},
    196600: {"ko": "상법및계약법", "en": "Commercial & Contract Law", "zh": "商法及契約法", "prof": "강편모"},
    196656: {"ko": "M&A전략", "en": "M&A Strategy", "zh": "M&A戰略", "prof": "김철중"},
    196622: {"ko": "국제거시금융론", "en": "International Macrofinance", "zh": "國際宏觀金融論", "prof": "이창민"},
}

# PPT 색상 테마 (한양대 느낌)
THEME = {
    "primary": RGBColor(0, 51, 102),      # 진한 남색
    "secondary": RGBColor(0, 102, 153),    # 파란색
    "accent": RGBColor(204, 0, 0),         # 빨간색 강조
    "bg_dark": RGBColor(0, 51, 102),       # 어두운 배경
    "bg_light": RGBColor(240, 245, 250),   # 밝은 배경
    "text_dark": RGBColor(33, 33, 33),     # 본문 텍스트
    "text_light": RGBColor(255, 255, 255), # 밝은 텍스트
    "text_sub": RGBColor(100, 100, 100),   # 부제 텍스트
}


def repair_truncated_json(text):
    """잘린 JSON을 복구 시도"""
    # 열린 brackets/braces 카운트
    brace_count = 0
    bracket_count = 0
    in_string = False
    escape_next = False

    for ch in text:
        if escape_next:
            escape_next = False
            continue
        if ch == '\\':
            escape_next = True
            continue
        if ch == '"' and not escape_next:
            in_string = not in_string
            continue
        if in_string:
            continue
        if ch == '{':
            brace_count += 1
        elif ch == '}':
            brace_count -= 1
        elif ch == '[':
            bracket_count += 1
        elif ch == ']':
            bracket_count -= 1

    if brace_count == 0 and bracket_count == 0:
        return text  # 이미 정상

    # 잘린 문자열 닫기 (열린 따옴표)
    stripped = text.rstrip()
    if stripped.endswith('\\'):
        stripped = stripped[:-1]

    # 마지막 완전한 항목까지 잘라내기
    # 불완전한 마지막 value 제거
    last_complete = stripped
    for pattern in [
        r',\s*"[^"]*":\s*"[^"]*$',     # 잘린 string value
        r',\s*"[^"]*":\s*$',             # key만 있고 value 없음
        r',\s*"[^"]*$',                  # 잘린 key
        r',\s*\{[^}]*$',                 # 잘린 object
        r',\s*"[^"]*":\s*\[[^\]]*$',     # 잘린 array value
    ]:
        match = re.search(pattern, last_complete)
        if match:
            last_complete = last_complete[:match.start()]
            break

    # 닫는 brackets/braces 추가
    result = last_complete
    brace_count = 0
    bracket_count = 0
    in_string = False
    escape_next = False
    for ch in result:
        if escape_next:
            escape_next = False
            continue
        if ch == '\\':
            escape_next = True
            continue
        if ch == '"':
            in_string = not in_string
            continue
        if in_string:
            continue
        if ch == '{':
            brace_count += 1
        elif ch == '}':
            brace_count -= 1
        elif ch == '[':
            bracket_count += 1
        elif ch == ']':
            bracket_count -= 1

    result += ']' * bracket_count + '}' * brace_count
    return result


def generate_ppt_content(topic, course_id=None, lang="ko", num_slides=18):
    """Claude AI로 PPT 내용 생성 (2-pass: 많은 슬라이드용)"""
    client = anthropic.Anthropic(api_key=API_KEY, timeout=180.0)

    course_info = ""
    if course_id and course_id in COURSE_MAP:
        c = COURSE_MAP[course_id]
        course_info = f"\n과목: {c['ko']} ({c['en']}) / 교수: {c['prof']}"

    lang_map = {"ko": "한국어", "en": "English", "zh": "中文繁體"}
    target_lang = lang_map.get(lang, "한국어")

    # 많은 슬라이드는 2번에 나눠 생성
    if num_slides > 12:
        first_half = num_slides // 2
        second_half = num_slides - first_half
        print(f"  [INFO] {num_slides}장 → 2회 생성 ({first_half}장 + {second_half}장)")

        part1 = _call_claude_for_slides(
            client, topic, course_info, target_lang,
            first_half, "Part 1: 서론~본론",
            include_meta=True
        )
        part2 = _call_claude_for_slides(
            client, topic, course_info, target_lang,
            second_half, "Part 2: 사례분석~결론~참고자료",
            include_meta=False,
            existing_titles=[s["title"] for s in part1.get("slides", [])]
        )

        # 합치기
        part1["slides"].extend(part2.get("slides", []))
        part1["key_terms"] = part1.get("key_terms", []) + part2.get("key_terms", [])
        return part1
    else:
        return _call_claude_for_slides(
            client, topic, course_info, target_lang,
            num_slides, "전체", include_meta=True
        )


def _call_claude_for_slides(client, topic, course_info, target_lang,
                             num_slides, part_desc, include_meta=True,
                             existing_titles=None):
    """Claude API 호출하여 슬라이드 JSON 생성"""

    existing_note = ""
    if existing_titles:
        titles_str = ", ".join(existing_titles)
        existing_note = f"\n\n이미 생성된 슬라이드: [{titles_str}]\n위 내용과 중복되지 않게 이어서 작성하세요."

    meta_fields = ""
    if include_meta:
        meta_fields = '"title": "PPT 제목",\n  "subtitle": "부제목",\n  '

    prompt = f"""MBA 수업 발표용 PPT 내용을 JSON으로 생성하세요.

주제: {topic}{course_info}
언어: {target_lang}
슬라이드 수: 정확히 {num_slides}장 ({part_desc})
{existing_note}

**중요: JSON만 출력하세요. 다른 설명 없이 순수 JSON만!**

슬라이드 type 종류 (반드시 2개 이상의 차트/테이블 포함!):
- "content": 일반 불릿 슬라이드
- "case": 사례 연구 (company 필드 필수)
- "chart": 차트/도표 (bar/pie/line 등)
- "table": 비교표/정리표
- "references": 참고자료 (마지막)

JSON 구조:
{{{{{meta_fields}"slides": [
    {{"title": "제목", "type": "content", "bullets": ["포인트1", "포인트2", "포인트3", "포인트4"], "notes": "발표자 노트"}},
    {{"title": "시장 규모 추이", "type": "chart", "chart_type": "bar", "categories": ["2022", "2023", "2024", "2025"], "series": [{{"name": "한국 M&A (조원)", "values": [45, 52, 68, 75]}}, {{"name": "글로벌 (조달러)", "values": [3.2, 2.8, 3.5, 3.9]}}], "bullets": ["핵심 인사이트1", "인사이트2"], "source": "출처", "notes": "발표자 노트"}},
    {{"title": "유형별 비율", "type": "chart", "chart_type": "pie", "categories": ["수평적", "수직적", "혼합형", "역합병"], "series": [{{"name": "비율", "values": [35, 28, 25, 12]}}], "bullets": ["설명1"], "source": "출처", "notes": "노트"}},
    {{"title": "추이 분석", "type": "chart", "chart_type": "line", "categories": ["Q1", "Q2", "Q3", "Q4"], "series": [{{"name": "거래건수", "values": [120, 145, 165, 180]}}], "bullets": ["트렌드 설명"], "source": "출처", "notes": "노트"}},
    {{"title": "비교 분석", "type": "table", "headers": ["구분", "장점", "단점", "사례"], "rows": [["수평적 M&A", "시장점유율 확대", "독점규제", "디즈니-폭스"], ["수직적 M&A", "공급망 통합", "유연성 저하", "아마존-홀푸드"]], "notes": "노트"}},
    {{"title": "사례", "type": "case", "company": "기업명", "bullets": ["사례1", "사례2", "사례3"], "notes": "노트"}},
    {{"title": "참고자료", "type": "references", "bullets": ["출처1", "출처2"], "notes": ""}}
  ],
  "key_terms": [{{"term": "용어", "ko": "한국어", "en": "English", "zh": "中文"}}]
}}}}

요구사항:
1. MBA 수준의 전문적이고 상세한 내용
2. content/case 슬라이드 bullets는 3-5개, 각 bullet은 1-2문장
3. **중요** chart 슬라이드 bullets는 2-3개, 각 bullet은 반드시 15자 이내로 짧게! (차트 옆 좁은 공간)
4. 차트 최소 2개 포함 (실제 통계 데이터 기반, 수치는 현실적으로!)
5. 테이블 최소 1개 포함 (비교분석 또는 정리표)
6. 실제 기업 사례 2-3개 포함 (type: "case")
7. 관련 논문/기사 참고자료 최소 5개 (마지막 슬라이드, type: "references")
8. 핵심 용어 한/영/중 3개 언어 병기 (최소 5개)
9. 각 슬라이드에 발표자 노트 포함 (2-3문장, 상세 설명은 여기에!)
10. 목차 슬라이드 포함 (type: "content", title: "목차 / Table of Contents")"""

    response = client.messages.create(
        model="claude-sonnet-4-5-20250929",
        max_tokens=16384,
        messages=[{"role": "user", "content": prompt}],
    )

    text = response.content[0].text.strip()
    return _parse_json_response(text, topic)


def _parse_json_response(text, fallback_topic):
    """JSON 응답 파싱 (복구 포함)"""
    # 1) ```json ... ``` 블록 추출
    json_match = re.search(r'```json\s*(.*?)\s*```', text, re.DOTALL)
    if json_match:
        raw = json_match.group(1)
    else:
        # 2) { 로 시작하는 JSON 찾기
        brace_start = text.find('{')
        if brace_start >= 0:
            raw = text[brace_start:]
        else:
            raw = text

    # 3) 파싱 시도
    try:
        return json.loads(raw)
    except json.JSONDecodeError as e:
        print(f"  [WARN] JSON 파싱 실패: {e}")
        print(f"  [INFO] 잘린 JSON 복구 시도...")

    # 4) 복구 시도
    try:
        repaired = repair_truncated_json(raw)
        result = json.loads(repaired)
        print(f"  [OK] JSON 복구 성공! (슬라이드 {len(result.get('slides', []))}장)")
        return result
    except json.JSONDecodeError:
        pass

    # 5) 최후 수단: 줄 단위로 잘라가며 파싱
    lines = raw.split('\n')
    for i in range(len(lines), 0, -1):
        chunk = '\n'.join(lines[:i])
        try:
            repaired = repair_truncated_json(chunk)
            result = json.loads(repaired)
            if result.get("slides"):
                print(f"  [OK] 부분 복구 성공! (슬라이드 {len(result.get('slides', []))}장)")
                return result
        except (json.JSONDecodeError, Exception):
            continue

    print("[!] JSON 복구 실패, 기본 구조로 생성")
    return {
        "title": fallback_topic,
        "subtitle": "한양대 MBA",
        "slides": [
            {"title": fallback_topic, "type": "content",
             "bullets": [text[:300]], "notes": ""}
        ],
        "key_terms": []
    }


def set_slide_bg(slide, color):
    """슬라이드 배경색 설정"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=None, bold=False, alignment=PP_ALIGN.LEFT):
    """텍스트 박스 추가"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.alignment = alignment
    if color:
        p.font.color.rgb = color
    return tf


def create_title_slide(prs, data, course_id=None):
    """표지 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, THEME["bg_dark"])

    # 제목
    add_text_box(slide, Inches(1), Inches(1.5), Inches(8), Inches(1.5),
                 data["title"], font_size=36, color=THEME["text_light"],
                 bold=True, alignment=PP_ALIGN.CENTER)

    # 부제
    add_text_box(slide, Inches(1), Inches(3.2), Inches(8), Inches(0.8),
                 data.get("subtitle", ""), font_size=20,
                 color=RGBColor(180, 200, 220), alignment=PP_ALIGN.CENTER)

    # 과목 + 날짜
    if course_id and course_id in COURSE_MAP:
        c = COURSE_MAP[course_id]
        info = f"{c['ko']} | {c['prof']} 교수"
    else:
        info = "한양대학교 MBA 경영전문대학원"

    add_text_box(slide, Inches(1), Inches(4.5), Inches(8), Inches(0.5),
                 info, font_size=14, color=RGBColor(150, 170, 190),
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(5.0), Inches(8), Inches(0.5),
                 datetime.now().strftime("%Y년 %m월 %d일"), font_size=12,
                 color=RGBColor(130, 150, 170), alignment=PP_ALIGN.CENTER)


def create_content_slide(prs, slide_data, slide_num):
    """내용 슬라이드 - 하나의 텍스트 프레임으로 겹침 방지"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, THEME["bg_light"])

    # 상단 색상 바
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(0.08)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = THEME["primary"]
    shape.line.fill.background()

    # 슬라이드 번호
    add_text_box(slide, Inches(0.5), Inches(0.2), Inches(0.5), Inches(0.4),
                 f"{slide_num:02d}", font_size=11, color=THEME["text_sub"])

    # 제목
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.7),
                 slide_data["title"], font_size=28, color=THEME["primary"],
                 bold=True)

    # 구분선
    shape = slide.shapes.add_shape(
        1, Inches(0.5), Inches(1.2), Inches(2), Inches(0.03)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = THEME["accent"]
    shape.line.fill.background()

    # 불릿 포인트 - 하나의 텍스트 프레임에 여러 문단 (겹침 방지!)
    bullets = slide_data.get("bullets", [])
    font_sz = 16 if len(bullets) <= 4 else 14 if len(bullets) <= 6 else 12
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(4.0))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {bullet}"
        p.font.size = Pt(font_sz)
        p.font.color.rgb = THEME["text_dark"]
        p.space_before = Pt(6)
        p.space_after = Pt(4)
        # 불릿 설정
        p.level = 0
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        buChar = pPr.makeelement('{http://schemas.openxmlformats.org/drawingml/2006/main}buChar', {'char': '●'})
        # 기존 불릿 제거 후 추가
        for existing in pPr.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}buChar'):
            pPr.remove(existing)
        for existing in pPr.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}buNone'):
            pPr.remove(existing)
        pPr.append(buChar)
        # 불릿 색상
        buClr = pPr.makeelement('{http://schemas.openxmlformats.org/drawingml/2006/main}buClr', {})
        srgbClr = buClr.makeelement('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr',
                                     {'val': '006699'})
        buClr.append(srgbClr)
        for existing in pPr.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}buClr'):
            pPr.remove(existing)
        pPr.append(buClr)
        # 불릿 크기
        buSzPct = pPr.makeelement('{http://schemas.openxmlformats.org/drawingml/2006/main}buSzPct',
                                   {'val': '70000'})
        for existing in pPr.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}buSzPct'):
            pPr.remove(existing)
        pPr.append(buSzPct)

    # 발표자 노트
    notes = slide_data.get("notes", "")
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes


def create_case_slide(prs, slide_data, slide_num):
    """사례 연구 슬라이드 - 하나의 텍스트 프레임으로 겹침 방지"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(255, 255, 255))

    # 상단 바
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(0.08)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = THEME["accent"]
    shape.line.fill.background()

    # 헤더 + 제목 + company를 하나의 텍스트 프레임으로
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1.8))
    tf = txBox.text_frame
    tf.word_wrap = True

    # CASE STUDY 라벨
    p = tf.paragraphs[0]
    p.text = "CASE STUDY"
    p.font.size = Pt(11)
    p.font.color.rgb = THEME["accent"]
    p.font.bold = True
    p.space_after = Pt(2)

    # 제목
    p2 = tf.add_paragraph()
    p2.text = slide_data["title"]
    p2.font.size = Pt(24)
    p2.font.color.rgb = THEME["primary"]
    p2.font.bold = True
    p2.space_after = Pt(2)

    # 기업명
    company = slide_data.get("company", "")
    if company:
        p3 = tf.add_paragraph()
        p3.text = company
        p3.font.size = Pt(14)
        p3.font.color.rgb = THEME["secondary"]
        p3.font.bold = True
        p3.space_after = Pt(4)

    # 불릿 - 별도 텍스트 프레임 (제목 아래 충분한 공간)
    bullets = slide_data.get("bullets", [])
    font_sz = 14 if len(bullets) <= 5 else 12
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(3.2))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = f"▸ {bullet}"
        p.font.size = Pt(font_sz)
        p.font.color.rgb = THEME["text_dark"]
        p.space_before = Pt(4)
        p.space_after = Pt(4)

    notes = slide_data.get("notes", "")
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def create_terms_slide(prs, key_terms):
    """핵심 용어 슬라이드 (3개 언어)"""
    if not key_terms:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, THEME["bg_light"])

    # 상단 바
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(0.08)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = THEME["primary"]
    shape.line.fill.background()

    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.7),
                 "Key Terms / 핵심 용어 / 關鍵術語", font_size=24,
                 color=THEME["primary"], bold=True)

    # 테이블 형식으로 용어 표시
    y_pos = 1.5
    # 헤더
    add_text_box(slide, Inches(0.5), Inches(1.2), Inches(3), Inches(0.4),
                 "한국어", font_size=13, color=THEME["secondary"], bold=True)
    add_text_box(slide, Inches(3.5), Inches(1.2), Inches(3), Inches(0.4),
                 "English", font_size=13, color=THEME["secondary"], bold=True)
    add_text_box(slide, Inches(6.5), Inches(1.2), Inches(3), Inches(0.4),
                 "中文", font_size=13, color=THEME["secondary"], bold=True)

    for term in key_terms[:8]:
        add_text_box(slide, Inches(0.5), Inches(y_pos), Inches(3), Inches(0.4),
                     term.get("ko", ""), font_size=14, color=THEME["text_dark"])
        add_text_box(slide, Inches(3.5), Inches(y_pos), Inches(3), Inches(0.4),
                     term.get("en", ""), font_size=14, color=THEME["text_dark"])
        add_text_box(slide, Inches(6.5), Inches(y_pos), Inches(3), Inches(0.4),
                     term.get("zh", ""), font_size=14, color=THEME["text_dark"])
        y_pos += 0.45


def create_references_slide(prs, slide_data):
    """참고자료 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(255, 255, 255))

    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(0.08)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = THEME["primary"]
    shape.line.fill.background()

    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.7),
                 "References / 참고 자료", font_size=24,
                 color=THEME["primary"], bold=True)

    # 참고자료 - 하나의 텍스트 프레임 (겹침 방지)
    bullets = slide_data.get("bullets", [])
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(4.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, ref in enumerate(bullets, 1):
        if i == 1:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"[{i}] {ref}"
        p.font.size = Pt(11)
        p.font.color.rgb = THEME["text_sub"]
        p.space_before = Pt(3)
        p.space_after = Pt(3)


def create_chart_slide(prs, slide_data, slide_num):
    """차트 슬라이드 (bar, pie, line)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(255, 255, 255))

    # 상단 바
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(0.08)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = THEME["secondary"]
    shape.line.fill.background()

    # 슬라이드 번호
    add_text_box(slide, Inches(0.5), Inches(0.2), Inches(0.5), Inches(0.4),
                 f"{slide_num:02d}", font_size=11, color=THEME["text_sub"])

    # 제목
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.6),
                 slide_data["title"], font_size=24, color=THEME["primary"],
                 bold=True)

    # 차트 데이터
    chart_type_str = slide_data.get("chart_type", "bar")
    categories = slide_data.get("categories", ["A", "B", "C"])
    series_list = slide_data.get("series", [{"name": "데이터", "values": [10, 20, 30]}])

    chart_data = CategoryChartData()
    chart_data.categories = categories
    for s in series_list:
        chart_data.add_series(s.get("name", ""), s.get("values", []))

    # 차트 타입 매핑
    chart_type_map = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "bar_stacked": XL_CHART_TYPE.COLUMN_STACKED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
        "area": XL_CHART_TYPE.AREA,
    }
    xl_chart_type = chart_type_map.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

    # 차트 위치: 제목 아래, 왼쪽 2/3
    chart_left = Inches(0.5)
    chart_top = Inches(1.3)
    chart_width = Inches(6)
    chart_height = Inches(3.8)

    chart_frame = slide.shapes.add_chart(
        xl_chart_type, chart_left, chart_top, chart_width, chart_height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(9)

    # 차트 색상 적용
    if chart_type_str == "pie":
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.size = Pt(10)
        data_labels.show_percentage = True
        data_labels.show_category_name = True
    else:
        for idx, series in enumerate(chart.series):
            fill = series.format.fill
            fill.solid()
            colors = [THEME["primary"], THEME["secondary"], THEME["accent"],
                      RGBColor(76, 153, 0), RGBColor(255, 153, 0)]
            fill.fore_color.rgb = colors[idx % len(colors)]

    # 오른쪽에 요약 불릿 - 하나의 텍스트 프레임 (겹침 방지)
    bullets = slide_data.get("bullets", [])
    if bullets:
        txBox = slide.shapes.add_textbox(Inches(6.6), Inches(1.3), Inches(3.2), Inches(3.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets[:4]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"▸ {bullet}"
            p.font.size = Pt(10)
            p.font.color.rgb = THEME["text_dark"]
            p.space_before = Pt(3)
            p.space_after = Pt(3)

    # 출처 표시
    source = slide_data.get("source", "")
    if source:
        add_text_box(slide, Inches(0.5), Inches(5.1), Inches(9), Inches(0.3),
                     f"출처: {source}", font_size=8, color=THEME["text_sub"])

    notes = slide_data.get("notes", "")
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def create_table_slide(prs, slide_data, slide_num):
    """테이블/비교표 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(255, 255, 255))

    # 상단 바
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(0.08)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = THEME["primary"]
    shape.line.fill.background()

    # 슬라이드 번호
    add_text_box(slide, Inches(0.5), Inches(0.2), Inches(0.5), Inches(0.4),
                 f"{slide_num:02d}", font_size=11, color=THEME["text_sub"])

    # 제목
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.6),
                 slide_data["title"], font_size=24, color=THEME["primary"],
                 bold=True)

    # 테이블 데이터
    headers = slide_data.get("headers", ["항목", "내용"])
    rows_data = slide_data.get("rows", [["예시", "데이터"]])

    num_rows = len(rows_data) + 1  # +1 for header
    num_cols = len(headers)

    # 테이블 위치/크기 자동 조정
    table_width = min(9.0, num_cols * 2.2)
    table_left = Inches((10 - table_width) / 2)  # 중앙 정렬
    row_height = min(0.5, 3.5 / max(num_rows, 1))

    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        table_left, Inches(1.3),
        Inches(table_width), Inches(row_height * num_rows + 0.3)
    )
    table = table_shape.table

    # 헤더 스타일
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = THEME["text_light"]
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = THEME["primary"]

    # 데이터 행
    for row_idx, row in enumerate(rows_data):
        for col_idx, value in enumerate(row):
            if col_idx >= num_cols:
                break
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.color.rgb = THEME["text_dark"]
            # 줄무늬 배경
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(245, 248, 252)

    notes = slide_data.get("notes", "")
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def create_end_slide(prs):
    """마지막 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, THEME["bg_dark"])

    add_text_box(slide, Inches(1), Inches(2.5), Inches(8), Inches(1),
                 "Thank You", font_size=40, color=THEME["text_light"],
                 bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(3.8), Inches(8), Inches(0.5),
                 "한양대학교 MBA 경영전문대학원", font_size=14,
                 color=RGBColor(150, 170, 190), alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(4.5), Inches(8), Inches(0.4),
                 "Generated by Smart Auto Study Bot", font_size=10,
                 color=RGBColor(100, 120, 140), alignment=PP_ALIGN.CENTER)


def build_ppt(data, course_id=None, output_path=None):
    """PPT 파일 생성"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    # 1. 표지
    create_title_slide(prs, data, course_id)

    # 2. 내용 슬라이드
    slide_num = 1
    for slide_data in data.get("slides", []):
        slide_type = slide_data.get("type", "content")

        if slide_type == "case":
            create_case_slide(prs, slide_data, slide_num)
        elif slide_type == "references":
            create_references_slide(prs, slide_data)
        elif slide_type in ("chart", "chart_bar", "chart_pie", "chart_line", "chart_area"):
            create_chart_slide(prs, slide_data, slide_num)
        elif slide_type == "table":
            create_table_slide(prs, slide_data, slide_num)
        else:
            create_content_slide(prs, slide_data, slide_num)
        slide_num += 1

    # 3. 핵심 용어 (3개 언어)
    create_terms_slide(prs, data.get("key_terms", []))

    # 4. 마지막 슬라이드
    create_end_slide(prs)

    # 저장
    if not output_path:
        safe_title = data["title"].replace(" ", "_").replace("/", "-")[:30]
        timestamp = datetime.now().strftime("%Y%m%d_%H%M")
        output_path = PROJECT_ROOT / "data" / "study_notes" / f"PPT_{safe_title}_{timestamp}.pptx"

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    return output_path


def main():
    parser = argparse.ArgumentParser(description="Smart Auto PPT Generator")
    parser.add_argument("--topic", required=True, help="PPT 주제")
    parser.add_argument("--course", type=int, help="Course ID")
    parser.add_argument("--lang", default="ko", choices=["ko", "en", "zh"], help="언어")
    parser.add_argument("--slides", type=int, default=18, help="슬라이드 수 (MBA 기본 18장)")
    parser.add_argument("--output", help="출력 파일 경로")
    parser.add_argument("--upload", action="store_true", help="Drive 업로드")
    parser.add_argument("--email", action="store_true", help="이메일로 발송")
    parser.add_argument("--whatsapp", action="store_true", help="WhatsApp 알림")
    parser.add_argument("--send", action="store_true", help="이메일 + WhatsApp 모두 발송")
    args = parser.parse_args()

    if not API_KEY:
        print("[!] ANTHROPIC_API_KEY 필요!")
        sys.exit(1)

    print(f"{'='*50}")
    print(f"  Smart Auto PPT Generator")
    print(f"  주제: {args.topic}")
    if args.course and args.course in COURSE_MAP:
        print(f"  과목: {COURSE_MAP[args.course]['ko']}")
    print(f"  언어: {args.lang} | 슬라이드: {args.slides}장")
    print(f"{'='*50}")

    # 1. AI 내용 생성
    print("\n[1/3] AI 내용 생성 중...")
    data = generate_ppt_content(args.topic, args.course, args.lang, args.slides)
    print(f"  제목: {data.get('title', args.topic)}")
    print(f"  슬라이드: {len(data.get('slides', []))}장")

    # 2. PPT 파일 생성
    print("\n[2/3] PPT 파일 생성 중...")
    output_path = build_ppt(data, args.course, args.output)
    print(f"  저장: {output_path}")
    print(f"  크기: {output_path.stat().st_size / 1024:.1f} KB")

    # 3. Drive 업로드 (선택)
    if args.upload:
        print("\n[3/3] Google Drive 업로드 중...")
        try:
            sys.path.insert(0, str(PROJECT_ROOT / "scrapers"))
            from drive_uploader import get_drive_service, find_or_create_folder, upload_file

            service = get_drive_service()
            DRIVE_ROOT = "1RGGiiz_DKf5ZcUNk7baJis95oAcDJ5Em"

            with open(PROJECT_ROOT / "config" / "courses.json", "r", encoding="utf-8") as f:
                config = json.load(f)

            semester_folder = find_or_create_folder(service, f"한양MBA_{config['semester']}", DRIVE_ROOT)

            if args.course:
                course_name = COURSE_MAP.get(args.course, {}).get("ko", "기타")
                course_folder = find_or_create_folder(service, course_name, semester_folder)
                notes_folder = find_or_create_folder(service, "학습노트", course_folder)
                file_id = upload_file(service, output_path, notes_folder)
            else:
                file_id = upload_file(service, output_path, semester_folder)

            drive_url = f"https://drive.google.com/file/d/{file_id}/view"
            print(f"  Drive URL: {drive_url}")
        except Exception as e:
            print(f"  [!] 업로드 실패: {e}")

    # 4. 이메일 발송
    if args.email or args.send:
        print("\n[EMAIL] 이메일 발송 중...")
        try:
            from email_sender import send_email_with_attachment
            course_name = COURSE_MAP.get(args.course, {}).get("ko", "") if args.course else ""
            subject = f"[Smart Auto] {course_name} - {data.get('title', args.topic)}"
            send_email_with_attachment(output_path, subject=subject)
        except Exception as e:
            print(f"  [!] 이메일 실패: {e}")

    # 5. WhatsApp 알림
    if args.whatsapp or args.send:
        print("\n[WHATSAPP] WhatsApp 알림 발송 중...")
        try:
            import subprocess
            target = os.getenv("WHATSAPP_TARGET", "")
            course_name = COURSE_MAP.get(args.course, {}).get("ko", "") if args.course else ""
            msg = (
                f"PPT 생성 완료!\n"
                f"제목: {data.get('title', args.topic)}\n"
                f"과목: {course_name}\n"
                f"슬라이드: {len(data.get('slides', []))}장\n"
                f"파일: {output_path.name}\n"
                f"이메일로도 보냈어요!"
            )
            subprocess.run(
                ["wsl", "bash", "-c",
                 f'npx openclaw message send --channel whatsapp '
                 f'--target "{target}" --message "{msg}" --json'],
                capture_output=True, timeout=30,
            )
            print("  WhatsApp 알림 발송 완료!")
        except Exception as e:
            print(f"  [!] WhatsApp 실패: {e}")

    print(f"\n{'='*50}")
    print(f"[DONE] PPT 생성 완료!")
    print(f"  파일: {output_path.name}")

    # JSON 결과 출력 (에이전트 파싱용)
    result = {
        "status": "success",
        "file": str(output_path),
        "file_name": output_path.name,
        "title": data.get("title", args.topic),
        "slides": len(data.get("slides", [])),
        "terms": len(data.get("key_terms", [])),
    }
    print(json.dumps(result, ensure_ascii=False))


if __name__ == "__main__":
    main()
